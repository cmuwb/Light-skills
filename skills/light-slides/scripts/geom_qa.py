# -*- coding: utf-8 -*-
"""geom_qa.py — PPT 确定性几何 QA（Light m16）。

补上审查指出的名实不符（claim_vs_code_gap #1/#2）：SKILL/thumbnail 反复宣称
"缩略图足以抓重叠/溢出/对齐错位"，但此前**没有任何代码真的计算重叠/溢出** ——
只把几何画成 PIL 示意图交人眼看。本脚本读每页 shape 的 left/top/width/height
(EMU)，调用共享地基契约 `_shared/visual_qa.detect_geometry_issues` 做**确定性**
重叠/溢出画布/超边距/对齐偏差检测，结果可并入 pptx_eval 的 layout 维度，
a08 闸门把"零重叠零溢出"设为硬通过条件。

可选渲染回看：有 LibreOffice/PyMuPDF 时可渲染 PNG 喂回多模态 Opus 按
visual_qa_rubric() 打分（render_then_review 协议见 visual_qa 模块）；无渲染器时
降级为仅几何检测并显式标注"未做像素级回看"，不静默假成功。

python-pptx 缺失时降级（CLI 报缺依赖；selftest 用合成 shape 验算法）。
`python geom_qa.py deck.pptx [--json]` / `python geom_qa.py --selftest`。
"""
from __future__ import annotations
import argparse
import json
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "..", "..", "_shared"))
try:
    from visual_qa import detect_geometry_issues, qa_report, visual_qa_rubric  # noqa: E402
    _HAS_VQA = True
except Exception:
    _HAS_VQA = False

EMU_PER_INCH = 914400
DEFAULT_W = int(13.333 * EMU_PER_INCH)   # 16:9 宽屏
DEFAULT_H = int(7.5 * EMU_PER_INCH)


def read_shapes(slide):
    """从 python-pptx slide 读出 AABB 列表(EMU)。供 geom_qa/thumbnail/pptx_eval 共用。"""
    shapes = []
    for i, sh in enumerate(slide.shapes):
        if sh.left is None or sh.top is None or not sh.width or not sh.height:
            continue
        kind = "text" if getattr(sh, "has_text_frame", False) and sh.has_text_frame else "shape"
        text = ""
        if kind == "text":
            try:
                text = sh.text_frame.text or ""
            except Exception:
                text = ""
        shapes.append({"id": f"sh{i}", "x": int(sh.left), "y": int(sh.top),
                       "w": int(sh.width), "h": int(sh.height), "kind": kind,
                       "text": text[:40]})
    return shapes


def qa_pptx(pptx_path, margins_frac=0.04):
    """对整个 deck 逐页做几何 QA。返回 {slides:[{idx, n_shapes, report}], summary}。"""
    if not _HAS_VQA:
        raise RuntimeError("visual_qa 共享契约不可用（_shared 未就位）")
    try:
        from pptx import Presentation
        from pptx.util import Emu  # noqa: F401
    except Exception as e:
        raise RuntimeError(f"python-pptx 未安装，无法读 .pptx（{e}）")
    prs = Presentation(pptx_path)
    W = int(prs.slide_width or DEFAULT_W)
    H = int(prs.slide_height or DEFAULT_H)
    margins = {"left": int(W * margins_frac), "right": int(W * margins_frac),
               "top": int(H * margins_frac), "bottom": int(H * margins_frac)}
    slides_out = []
    n_crit = n_imp = 0
    for idx, slide in enumerate(prs.slides, 1):
        shapes = read_shapes(slide)
        issues = detect_geometry_issues(shapes, (W, H), margins,
                                        align_tol=int(W * 0.005))
        rep = qa_report(issues, vlm_findings=None)  # 仅几何，未做像素回看
        n_crit += rep["counts"]["critical"]
        n_imp += rep["counts"]["important"]
        slides_out.append({"idx": idx, "n_shapes": len(shapes), "report": rep})
    return {"schema": "light.geom_qa.v1", "pptx": os.path.abspath(pptx_path),
            "n_slides": len(slides_out),
            "summary": {"critical": n_crit, "important": n_imp,
                        "hard_pass": n_crit == 0,
                        "note": "仅确定性几何检测；像素级视觉品味/对比度需 render_then_review 回看"},
            "slides": slides_out}


def _selftest() -> int:
    ok = True

    def check(cond, msg):
        nonlocal ok
        if not cond:
            ok = False
        print(f"  [{'PASS' if cond else 'FAIL'}] {msg}")

    print("geom_qa selftest")
    check(_HAS_VQA, "visual_qa 共享契约可导入")
    if not _HAS_VQA:
        print("SOME FAILED"); return 1

    W, H = DEFAULT_W, DEFAULT_H
    margins = {"left": int(W * 0.04), "right": int(W * 0.04),
               "top": int(H * 0.04), "bottom": int(H * 0.04)}

    # 1. 两个重叠文本框 → 检出 overlap
    overlap = [
        {"id": "t1", "x": int(W * 0.1), "y": int(H * 0.1), "w": int(W * 0.5), "h": int(H * 0.3), "kind": "text"},
        {"id": "t2", "x": int(W * 0.3), "y": int(H * 0.2), "w": int(W * 0.5), "h": int(H * 0.3), "kind": "text"},
    ]
    iss = detect_geometry_issues(overlap, (W, H), margins)
    check(any(i["type"] == "overlap" for i in iss), "检出 PPT 元素重叠")

    # 2. 溢出画布 → critical
    overflow = [{"id": "big", "x": int(W * 0.7), "y": int(H * 0.1), "w": int(W * 0.5), "h": int(H * 0.3)}]
    iss2 = detect_geometry_issues(overflow, (W, H), margins)
    check(any(i["type"] == "overflow_canvas" for i in iss2), "检出溢出画布")
    rep2 = qa_report(iss2)
    check(rep2["verdict"] == "fail" and not rep2["pixel_review_done"],
          "几何 critical→fail 且标注未做像素回看")

    # 3. 正常布局 → 零硬问题，hard_pass
    clean = [
        {"id": "title", "x": int(W * 0.06), "y": int(H * 0.06), "w": int(W * 0.88), "h": int(H * 0.15), "kind": "text"},
        {"id": "body", "x": int(W * 0.06), "y": int(H * 0.30), "w": int(W * 0.88), "h": int(H * 0.55), "kind": "text"},
    ]
    iss3 = detect_geometry_issues(clean, (W, H), margins)
    check(not any(i["type"] in ("overlap", "overflow_canvas") for i in iss3),
          "正常布局无重叠/溢出误报")

    # 4. rubric 可用（供渲染回看喂多模态 Opus）
    check("dimensions" in visual_qa_rubric(), "visual_qa_rubric 可用于渲染回看")

    print("ALL PASS" if ok else "SOME FAILED")
    return 0 if ok else 1


def main():
    ap = argparse.ArgumentParser(description="PPT 确定性几何 QA（重叠/溢出/超边距/对齐）")
    ap.add_argument("pptx", nargs="?", help=".pptx 路径")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--selftest", action="store_true")
    a = ap.parse_args()
    if a.selftest:
        sys.exit(_selftest())
    if not a.pptx:
        print(__doc__); return
    try:
        res = qa_pptx(a.pptx)
    except RuntimeError as e:
        print(f"[降级] {e}")
        sys.exit(2)
    if a.json:
        print(json.dumps(res, ensure_ascii=False, indent=2))
    else:
        s = res["summary"]
        print(f"{res['n_slides']} 页：critical={s['critical']} important={s['important']} "
              f"hard_pass={s['hard_pass']}")
        for sl in res["slides"]:
            if sl["report"]["counts"]["total"]:
                print(f"  页{sl['idx']}: {sl['report']['counts']}")
    sys.exit(0 if res["summary"]["hard_pass"] else 1)


if __name__ == "__main__":
    main()
