"""assemble_from_spec.py — imggen-enhanced 流水线 Stage D：按 deck_spec 重组装配可编辑 pptx。

读 deck_spec.yaml + assets_gen/（生图元素）+ figures/（m11 真数据图）→ python-pptx 产出
**可编辑** pptx：文本是原生文本框（可改字/字号/颜色）、表格是原生表格、图片按 bbox 摆放。
这是与"整页贴图假 PPT"的本质区别——成品里每个元素都是独立可编辑对象。

契约与字段语义见 references/imggen_pipeline.md（Stage A/D）与 templates/deck_spec.yaml。
坐标换算：left = bbox.x × 13.333in，top = bbox.y × 7.5in（16:9 宽屏 EMU）。
复用 patterns.md 的 fill_bg/add_text/rect 思路（本脚本内联等价实现，装到 ~/.claude 后无法跨技能 import）。

用法：
    python assemble_from_spec.py spec.yaml --assets assets_gen --figures figures -o out.pptx
    python assemble_from_spec.py --selftest          # 完全离线：内置最小 spec + PIL 占位图，产 2 页断言后清理

硬边界（与 imggen_pipeline.md 一致）：chart 只 add_picture 真数据图、绝不生图；文本一律文本框、不烤进图。
无对应主题字体时回退 Microsoft YaHei 并 warn，不静默错排。
"""
from __future__ import annotations

import sys

sys.stdout.reconfigure(encoding="utf-8")
import os
import argparse
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 16:9 宽屏画布（与 patterns.md 一致）
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
FONT_FALLBACK = "Microsoft YaHei"   # 中文字体缺失时回退

# 主题色板：优先从 assets/themes.py 取；取不到用内置兜底（保证离线 selftest 不依赖 import 路径）
_FALLBACK_THEME = {
    "bg": "FFFFFF", "surface": "F2F5F9", "primary": "1F4E79", "secondary": "2E75B6",
    "accent": "C00000", "text": "333333", "muted": "808080", "line": "D9D9D9",
}


def _load_theme_colors(theme_name: str) -> dict:
    """取主题 8 字段色板；themes.py 可 import 就用它，否则用兜底（离线稳健）。"""
    try:
        sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets"))
        from themes import get_theme  # type: ignore
        return dict(get_theme(theme_name)["COLORS"])
    except Exception:
        return dict(_FALLBACK_THEME)


def C(h: str) -> RGBColor:
    """hex(6 位无 #) -> RGBColor。"""
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _resolve_color(value, palette: dict) -> str:
    """color 字段可写成 hex 或调色板键名（如 'text'/'primary'）；解析成 hex。"""
    if not value:
        return palette.get("text", "333333")
    return palette.get(value, value)  # 是键名取键值，否则当 hex 原样返回


def rect(slide, x, y, w, h, color):
    """实心无边无影矩形（色块/色带/背景），坐标单位相对 0-1。"""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x * SLIDE_W_IN), Inches(y * SLIDE_H_IN),
                                Inches(w * SLIDE_W_IN), Inches(h * SLIDE_H_IN))
    sp.fill.solid(); sp.fill.fore_color.rgb = C(color)
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def fill_bg(slide, prs, color):
    """页背景：全幅矩形沉到底层（python-pptx 无原生页背景，patterns.md §0 同法）。"""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                prs.slide_width, prs.slide_height)
    sp.fill.solid(); sp.fill.fore_color.rgb = C(color)
    sp.line.fill.background(); sp.shadow.inherit = False
    slide.shapes._spTree.remove(sp._element)
    slide.shapes._spTree.insert(2, sp._element)
    return sp


def add_text(slide, bbox, lines, size, color, *, bold=False, font=FONT_FALLBACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """原生文本框，margin=0 才能与形状精确对齐。lines 可为 str 或 [str,...]（每条一段）。"""
    x, y, w, h = bbox
    tb = slide.shapes.add_textbox(Inches(x * SLIDE_W_IN), Inches(y * SLIDE_H_IN),
                                  Inches(w * SLIDE_W_IN), Inches(h * SLIDE_H_IN))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    items = lines if isinstance(lines, list) else [lines]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ("• " + item) if (isinstance(lines, list) and len(items) > 1) else item
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
        r.font.color.rgb = C(color)
    return tb


def place_element(slide, prs, el, palette, fonts, *, assets_dir, figures_dir, warnings):
    """按 type 分流单个元素。返回是否计入"文本元素"（QA 文本未栅格化检查用）。"""
    t = el.get("type")
    bbox = el.get("bbox", [0, 0, 1, 1])
    color = _resolve_color(el.get("color"), palette)

    if t == "title":
        add_text(slide, bbox, el.get("text", ""), el.get("font_size", 28), color,
                 bold=True, font=fonts["title"])
        return True
    if t == "body":
        content = el.get("bullets") or el.get("text", "")
        add_text(slide, bbox, content, el.get("font_size", 18), color, font=fonts["body"])
        return True
    if t == "table":
        _place_table(slide, el, palette, fonts, warnings)
        return False
    if t in ("chart", "icon", "illustration", "decor", "background"):
        # chart=真数据图(figures/)，其余=生图元素(assets_gen/)；统一 add_picture，绝不在此生数据图
        src = el.get("source") or el.get("data_ref")
        base = figures_dir if t == "chart" else assets_dir
        path = src if (src and os.path.isabs(src)) else os.path.join(base or ".", os.path.basename(src)) if src else None
        if path and os.path.exists(path):
            x, y, w, h = bbox
            slide.shapes.add_picture(path, Inches(x * SLIDE_W_IN), Inches(y * SLIDE_H_IN),
                                     Inches(w * SLIDE_W_IN), Inches(h * SLIDE_H_IN))
        else:
            # 资产缺失：画占位框 + warn，不静默跳过（与诚信纪律一致）
            rect(slide, *bbox, palette.get("line", "D9D9D9"))
            warnings.append(f"{el.get('id','?')}({t}) 资产缺失：{path or src}")
        return False
    warnings.append(f"{el.get('id','?')} 未知 type={t!r}，跳过")
    return False


def _place_table(slide, el, palette, fonts, warnings):
    """原生表格（add_table）：表头主色、斑马纹。data 取 el['rows'] 或退化为占位 2x2。"""
    rows = el.get("rows") or [["维度", "值"], ["占位", "—"]]
    x, y, w, h = el.get("bbox", [0.1, 0.5, 0.8, 0.3])
    nr, nc = len(rows), max(len(r) for r in rows)
    tbl = slide.shapes.add_table(nr, nc, Inches(x * SLIDE_W_IN), Inches(y * SLIDE_H_IN),
                                 Inches(w * SLIDE_W_IN), Inches(h * SLIDE_H_IN)).table
    for i, row in enumerate(rows):
        for j in range(nc):
            cell = tbl.cell(i, j)
            cell.text = str(row[j]) if j < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = C(palette["primary"] if i == 0 else
                                         (palette["surface"] if i % 2 else palette["bg"]))
            pr = cell.text_frame.paragraphs[0]
            pr.font.size = Pt(14); pr.font.name = fonts["body"]
            pr.font.color.rgb = C("FFFFFF" if i == 0 else palette["text"])
            pr.font.bold = (i == 0)


def assemble(spec: dict, out_path: str, *, assets_dir=None, figures_dir=None) -> dict:
    """按 spec 装配可编辑 pptx。返回统计 {pages, elements, text_elements, warnings}。"""
    deck = spec.get("deck", {})
    palette = _load_theme_colors(deck.get("theme", "academic"))
    palette.update({k: v.lstrip("#") for k, v in (deck.get("palette") or {}).items()})  # spec 覆盖
    fp = deck.get("font_pair") or {}
    fonts = {"title": fp.get("title", FONT_FALLBACK), "body": fp.get("body", FONT_FALLBACK)}
    warnings: list[str] = []
    # 字体回退提醒：非 ASCII 字体名落地后若目标机缺字体，PowerPoint 会回退，这里先 warn
    for role, name in fonts.items():
        if not name:
            fonts[role] = FONT_FALLBACK
            warnings.append(f"{role} 字体为空，回退 {FONT_FALLBACK}")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    n_elem = n_text = 0
    for page in spec.get("pages", []):
        slide = prs.slides.add_slide(blank)
        fill_bg(slide, prs, palette["bg"])
        # decor/background 先铺底，其余后摆（保证装饰不压正文）
        els = sorted(page.get("elements", []),
                     key=lambda e: 0 if e.get("type") in ("decor", "background") else 1)
        for el in els:
            is_text = place_element(slide, prs, el, palette, fonts,
                                    assets_dir=assets_dir, figures_dir=figures_dir, warnings=warnings)
            n_elem += 1
            n_text += int(is_text)
        if page.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(page["notes"])

    os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
    prs.save(out_path)
    return {"pages": len(spec.get("pages", [])), "elements": n_elem,
            "text_elements": n_text, "warnings": warnings, "out": out_path}


_MIN_SPEC = {
    "deck": {"title": "selftest", "aspect": "16:9", "theme": "tech",
             "style_anchor": "扁平测试风", "font_pair": {"title": "Microsoft YaHei", "body": "Microsoft YaHei"}},
    "pages": [
        {"page_id": "p01", "page_type": "cover", "action_title": "封面",
         "elements": [
             {"id": "e1", "type": "title", "bbox": [0.07, 0.3, 0.86, 0.18], "text": "标题", "color": "FFFFFF"},
             {"id": "e2", "type": "decor", "bbox": [0.0, 0.0, 1.0, 1.0], "desc": "角落装饰"}],
         "notes": "开场"},
        {"page_id": "p02", "page_type": "content", "action_title": "内容",
         "elements": [
             {"id": "e1", "type": "title", "bbox": [0.06, 0.06, 0.88, 0.12], "text": "结论句", "color": "text"},
             {"id": "e2", "type": "body", "bbox": [0.06, 0.24, 0.5, 0.6], "bullets": ["要点一", "要点二"]},
             {"id": "e3", "type": "icon", "bbox": [0.66, 0.3, 0.1, 0.16], "desc": "图标"},
             {"id": "e4", "type": "table", "bbox": [0.06, 0.7, 0.5, 0.2], "rows": [["A", "B"], ["1", "2"]]}],
         "notes": "讲要点"},
    ],
}


def _selftest() -> int:
    """完全离线：内置最小 spec + PIL 占位图，产 2 页 pptx 到临时目录、断言后清理，不留产物。"""
    from pptx import Presentation as _P
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from imagegen import generate as _gen  # 同目录，selftest 时可 import

    with tempfile.TemporaryDirectory() as td:
        assets = os.path.join(td, "assets_gen")
        # 为 p02 的 icon 造一张 mock 占位图（命名规则 page_elem_type，basename 必须匹配 source）
        _gen("图标占位", os.path.join(assets, "icon.png"), backend="mock", kind="icon")
        spec = {**_MIN_SPEC}
        spec["pages"][1]["elements"][2]["source"] = "icon.png"  # 指到 mock 图
        out = os.path.join(td, "out.pptx")
        stats = assemble(spec, out, assets_dir=assets, figures_dir=None)

        assert os.path.exists(out), "pptx 未产出"
        assert stats["pages"] == 2, stats
        assert stats["elements"] == 6, stats          # 2 + 4
        assert stats["text_elements"] == 3, stats     # title×2 + body×1（table 不计文本）
        # 真打开校验：2 页、文本框存在（证明文字是原生文本框，未栅格化——硬边界三）
        prs = _P(out)
        assert len(prs.slides) == 2
        texts = []
        for sl in prs.slides:
            for sh in sl.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    texts.append(sh.text_frame.text)
        joined = " ".join(texts)
        assert "结论句" in joined and "要点一" in joined, f"原生文本缺失：{texts}"
        # notes 写入校验
        assert prs.slides[0].notes_slide.notes_text_frame.text == "开场"
        del prs  # 释放句柄（Windows）
    print("[selftest] PASS assemble_from_spec（2 页可编辑 pptx + 原生文本框 + notes，全离线已清理）")
    return 0


def main():
    ap = argparse.ArgumentParser(description="按 deck_spec 重组装配可编辑 pptx（Stage D）")
    ap.add_argument("spec", nargs="?", help="deck_spec.yaml 路径")
    ap.add_argument("-o", "--out", default="deck_assembled.pptx", help="输出 pptx 路径")
    ap.add_argument("--assets", dest="assets_dir", help="生图元素目录（assets_gen/）")
    ap.add_argument("--figures", dest="figures_dir", help="真数据图目录（figures/，m11 出图）")
    ap.add_argument("--selftest", action="store_true", help="完全离线自测")
    args = ap.parse_args()

    if args.selftest:
        raise SystemExit(_selftest())

    if not args.spec:
        ap.error("需要提供 deck_spec.yaml 路径（或用 --selftest）")
    import yaml
    with open(args.spec, encoding="utf-8") as f:
        spec = yaml.safe_load(f)
    stats = assemble(spec, args.out, assets_dir=args.assets_dir, figures_dir=args.figures_dir)
    print(f"[ok] {stats['pages']} 页 / {stats['elements']} 元素 -> {stats['out']}")
    if stats["warnings"]:
        print("  warnings:")
        for w in stats["warnings"]:
            print(f"    ! {w}")


if __name__ == "__main__":
    main()
